import base64
import logging
from pathlib import Path
from typing import List
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader

logger = logging.getLogger(__name__)


def _make_single_doc(text: str, file_path: str) -> List[Document]:
    return [Document(page_content=text, metadata={"source": file_path, "page": 0})]


class DocumentLoaderError(Exception):
    """Base exception for document loading failures with a user-facing message."""

    def __init__(self, message: str, detail: str = ""):
        super().__init__(message)
        self.detail = detail


class UnsupportedFileTypeError(DocumentLoaderError):
    """File extension not in the supported types list."""

    def __init__(self, ext: str):
        from config import ALLOWED_EXTENSIONS
        super().__init__(
            f"不支持的文件类型: {ext}",
            f"支持的类型: {', '.join(sorted(ALLOWED_EXTENSIONS))}",
        )


class EmptyDocumentError(DocumentLoaderError):
    """Document has no extractable text."""

    def __init__(self, detail: str = ""):
        super().__init__(
            "文档内容为空或无法提取有效文本。PDF 可能是扫描版（图片），请使用带有文字层的文档。",
            detail,
        )


class DocumentParseError(DocumentLoaderError):
    """Document format is unsupported or corrupted."""

    def __init__(self, detail: str = ""):
        super().__init__(
            "无法解析文档，文件可能已损坏或格式不正确。",
            detail,
        )

_OCR_CLIENT = None


def _get_ocr_client() -> "OpenAI":
    global _OCR_CLIENT
    if _OCR_CLIENT is None:
        from openai import OpenAI
        from config import DASHSCOPE_API_KEY, BASE_URL
        _OCR_CLIENT = OpenAI(api_key=DASHSCOPE_API_KEY, base_url=BASE_URL)
    return _OCR_CLIENT


def _ocr_image(image_bytes: bytes) -> str:
    """使用 DashScope 多模态模型 (qwen-vl) 进行 OCR 识别"""
    client = _get_ocr_client()
    img_b64 = base64.b64encode(image_bytes).decode("utf-8")

    response = client.chat.completions.create(
        model="qwen-vl-max",
        messages=[{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
                {"type": "text", "text": "请完整提取并输出这张图片中的所有文字内容，保持原文格式和顺序，不要添加任何额外说明。"},
            ],
        }],
        temperature=0.1,
        max_tokens=4096,
    )
    return response.choices[0].message.content or ""


def load_pdf(file_path: str) -> List[Document]:
    loader = PyPDFLoader(file_path)
    docs = loader.load()
    total_text = "".join(d.page_content for d in docs).strip()
    if total_text:
        return docs

    logger.info("PyPDF 未提取到文本，尝试 PyMuPDF...")

    import fitz
    fitz_doc = fitz.open(file_path)
    try:
        pages_text = []
        for i in range(fitz_doc.page_count):
            text = fitz_doc[i].get_text().strip()
            if text:
                pages_text.append(text)
        fitz_text = "\n\n".join(pages_text).strip()

        if fitz_text:
            return _make_single_doc(fitz_text, file_path)

        logger.info("PyMuPDF 也未提取到文本，启动 OCR 识别（可能较慢）...")

        ocr_pages = []
        for i in range(fitz_doc.page_count):
            logger.info(f"  OCR 处理第 {i+1}/{fitz_doc.page_count} 页...")
            page = fitz_doc[i]
            pix = page.get_pixmap(dpi=300)
            img_bytes = pix.tobytes("png")
            text = _ocr_image(img_bytes)
            if text:
                ocr_pages.append(f"--- 第 {i+1} 页 ---\n{text}")

        ocr_text = "\n\n".join(ocr_pages).strip()
        if not ocr_text:
            raise EmptyDocumentError("OCR 未能识别出文字，文档可能为纯图片或手写内容。")

        return _make_single_doc(ocr_text, file_path)
    finally:
        fitz_doc.close()


def load_image(file_path: str) -> List[Document]:
    logger.info("正在对图片进行 OCR 识别...")
    with open(file_path, "rb") as f:
        img_bytes = f.read()
    text = _ocr_image(img_bytes)
    if not text.strip():
        raise EmptyDocumentError("OCR 未能从图片中识别出文字。")
    return _make_single_doc(text, file_path)


def _load_doc_via_com(file_path: str) -> str | None:
    """Extract text from .doc using Word COM automation (Windows, requires pywin32)."""
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return None

    word = None
    try:
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(file_path)
        try:
            text = doc.Content.Text
            return text.strip() or None
        finally:
            doc.Close()
    except Exception as e:
        logger.warning(f"Word COM 提取失败: {e}")
        return None
    finally:
        if word is not None:
            try:
                word.Quit()
            except Exception as e:
                logger.warning(f"Word COM 清理失败 (Quit): {e}")
        try:
            pythoncom.CoUninitialize()
        except Exception as e:
            logger.warning(f"Word COM 清理失败 (CoUninitialize): {e}")


_DOC_ENCODINGS = ("utf-16-le", "gbk", "utf-8", "latin-1")
_DOC_READABLE_RE = (
    r"[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef"   # CJK
    r"a-zA-Z0-9\s\.\,\;\:\!\?\(\)（）《》""''"     # ASCII words + common punctuation
    r"\u00a0-\u00ff"                                # Latin-1 supplement
    r"]+"
)


def _extract_text_from_doc_raw(file_path: str) -> str | None:
    """Last resort: extract readable text from .doc binary without any dependencies."""
    import re
    try:
        with open(file_path, "rb") as f:
            data = f.read()
    except OSError:
        return None

    text = None
    for encoding in _DOC_ENCODINGS:
        try:
            text = data.decode(encoding, errors="ignore")
            break
        except Exception:
            continue
    if text is None:
        return None

    chunks = re.findall(_DOC_READABLE_RE, text)
    result = "".join(c for c in chunks if len(c) > 3)
    return result if len(result) > 80 else None


def load_doc(file_path: str) -> List[Document]:
    """处理旧版 .doc 格式 — 四级降级策略"""
    import subprocess

    # Level 1: python-docx（部分 .doc 实际是 docx 格式）
    try:
        docs = load_docx(file_path)
        if docs and docs[0].page_content.strip():
            return docs
    except Exception:
        pass

    # Level 2: antiword（如果系统装了）
    try:
        result = subprocess.run(
            ["antiword", file_path], capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return _make_single_doc(result.stdout, file_path)
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # Level 3: Word COM 自动化（Windows，需 pywin32）
    text = _load_doc_via_com(file_path)
    if text:
        return _make_single_doc(text, file_path)

    # Level 4: 原始二进制提取（零依赖最后兜底）
    text = _extract_text_from_doc_raw(file_path)
    if text:
        return _make_single_doc(text, file_path)

    raise DocumentParseError(
        "无法解析 .doc 文件。请用 Word 打开后另存为 .docx 格式再上传，"
        "或执行 pip install pywin32 启用 Word COM 自动化解析。"
    )


def load_docx(file_path: str) -> List[Document]:
    from docx import Document as DocxDocument
    doc = DocxDocument(file_path)
    full_text = [para.text for para in doc.paragraphs if para.text.strip()]
    text = "\n\n".join(full_text)
    return _make_single_doc(text, file_path)


def load_pptx(file_path: str) -> List[Document]:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        if texts:
            slides.append(f"--- 幻灯片 {i+1} ---\n" + "\n".join(texts))
    text = "\n\n".join(slides)
    return _make_single_doc(text, file_path)


def load_xlsx(file_path: str) -> List[Document]:
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            sheets_text.append(f"--- 工作表: {sheet_name} ---\n" + "\n".join(rows))
    wb.close()
    text = "\n\n".join(sheets_text)
    return _make_single_doc(text, file_path)


def load_html(file_path: str) -> List[Document]:
    from bs4 import BeautifulSoup
    with open(file_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    text = "\n".join(lines)
    return _make_single_doc(text, file_path)


def load_text(file_path: str) -> List[Document]:
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()
    return _make_single_doc(text, file_path)


LOADERS = {
    ".pdf": load_pdf,
    ".doc": load_doc,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_xlsx,
    ".html": load_html,
    ".htm": load_html,
    ".png": load_image,
    ".jpg": load_image,
    ".jpeg": load_image,
    ".bmp": load_image,
    ".tiff": load_image,
    ".tif": load_image,
    ".txt": load_text,
    ".md": load_text,
    ".csv": load_text,
    ".json": load_text,
    ".xml": load_text,
    ".log": load_text,
}


def load_file(file_path: str) -> List[Document]:
    ext = Path(file_path).suffix.lower()
    if ext not in LOADERS:
        raise UnsupportedFileTypeError(ext)

    loader = LOADERS[ext]
    docs = loader(file_path)
    fname = Path(file_path).name
    return [
        Document(page_content=doc.page_content, metadata={**doc.metadata, "filename": fname, "file_type": ext})
        for doc in docs
    ]
